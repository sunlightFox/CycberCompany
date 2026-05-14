from __future__ import annotations

import io
import json
import re
from dataclasses import dataclass
from typing import TYPE_CHECKING, Any, cast

from brain import BrainRouteRequest, ModelRouter
from brain.adapters import CancelToken, ModelAdapterError, ModelChatRequest, OpenAICompatibleClient
from core_types import ErrorCode, RiskLevel, TaskArtifact, TraceSpanStatus, TraceSpanType
from trace_service import TraceService, redact

from app.core.errors import AppError
from app.db.repositories.brain_repo import BrainRepository
from app.schemas.tasks import ToolExecuteRequest
from app.services.artifacts import ArtifactStore
from app.services.audit import AuditEventService
from app.services.model_routing import ModelRoutingService
from app.services.office_productivity import LOCAL_OFFICE_PROVIDER, office_governance_result
from app.services.secrets import SecretStore

DOCX_IMPORT_ERROR: ImportError | None
try:
    from docx import Document as _Document
    from docx.shared import Pt as _DocxPt
except ImportError as exc:  # pragma: no cover - exercised only when env is incomplete
    DOCX_IMPORT_ERROR = exc
    Document: Any = None
    Pt: Any = None
else:
    DOCX_IMPORT_ERROR = None
    Document = cast(Any, _Document)
    Pt = cast(Any, _DocxPt)

XLSX_IMPORT_ERROR: ImportError | None
try:
    from openpyxl import Workbook as _Workbook
    from openpyxl import load_workbook as _load_workbook
    from openpyxl.chart import BarChart as _BarChart
    from openpyxl.chart import Reference as _Reference
    from openpyxl.styles import Font as _Font
except ImportError as exc:  # pragma: no cover - exercised only when env is incomplete
    XLSX_IMPORT_ERROR = exc
    Workbook: Any = None
    load_workbook: Any = None
    BarChart: Any = None
    Reference: Any = None
    Font: Any = None
else:
    XLSX_IMPORT_ERROR = None
    Workbook = cast(Any, _Workbook)
    load_workbook = cast(Any, _load_workbook)
    BarChart = cast(Any, _BarChart)
    Reference = cast(Any, _Reference)
    Font = cast(Any, _Font)

PPTX_IMPORT_ERROR: ImportError | None
try:
    from pptx import Presentation as _Presentation
    from pptx.util import Inches as _Inches
    from pptx.util import Pt as _PptPt
except ImportError as exc:  # pragma: no cover - exercised only when env is incomplete
    PPTX_IMPORT_ERROR = exc
    Presentation: Any = None
    Inches: Any = None
    PptPt: Any = None
else:
    PPTX_IMPORT_ERROR = None
    Presentation = cast(Any, _Presentation)
    Inches = cast(Any, _Inches)
    PptPt = cast(Any, _PptPt)

if TYPE_CHECKING:
    from docx.document import Document as DocxDocument
    from openpyxl import Workbook as OpenpyxlWorkbook
    from pptx.presentation import Presentation as PptxPresentation
else:
    DocxDocument = object
    OpenpyxlWorkbook = object
    PptxPresentation = object

DOCX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
XLSX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


@dataclass(frozen=True)
class OfficeComposedContent:
    title: str
    summary: str
    model_used: bool
    strategy: str
    brain_id: str | None = None
    finish_reason: str | None = None


class OfficeContentComposer:
    def __init__(
        self,
        *,
        brain_repo: BrainRepository | None = None,
        model_routing_service: ModelRoutingService | None = None,
        secret_store: SecretStore | None = None,
        trace_service: TraceService | None = None,
        audit_service: AuditEventService | None = None,
    ) -> None:
        self._brains = brain_repo
        self._routing = model_routing_service
        self._secrets = secret_store
        self._trace = trace_service
        self._audit = audit_service
        self._router = ModelRouter()

    async def compose(
        self,
        kind: str,
        args: dict[str, Any],
        *,
        trace_id: str | None,
        task_id: str | None,
        member_id: str | None,
    ) -> OfficeComposedContent:
        fallback = self._fallback(kind, args)
        if not self._model_available(args, trace_id):
            return fallback
        assert trace_id is not None
        model_result = await self._compose_with_model(
            kind,
            args,
            fallback=fallback,
            trace_id=trace_id,
            task_id=task_id,
            member_id=member_id,
        )
        return model_result or fallback

    def _fallback(self, kind: str, args: dict[str, Any]) -> OfficeComposedContent:
        title = _text(args.get("title") or _default_title(args), _fallback_title(kind))
        summary = _text(
            args.get("summary") or args.get("subtitle") or args.get("content") or args.get("goal"),
            "",
        )
        return OfficeComposedContent(
            title=title,
            summary=summary,
            model_used=False,
            strategy="deterministic_template",
        )

    def _model_available(self, args: dict[str, Any], trace_id: str | None) -> bool:
        if not trace_id:
            return False
        if args.get("disable_model") or args.get("model_used") is False:
            return False
        return all(
            service is not None
            for service in (self._brains, self._routing, self._secrets, self._trace)
        )

    async def _compose_with_model(
        self,
        kind: str,
        args: dict[str, Any],
        *,
        fallback: OfficeComposedContent,
        trace_id: str,
        task_id: str | None,
        member_id: str | None,
    ) -> OfficeComposedContent | None:
        assert self._brains is not None
        assert self._routing is not None
        assert self._secrets is not None
        assert self._trace is not None
        config = await self._routing.get_config()
        available_brains = await self._brains.list_routable_brains()
        route_selection = self._router.select_route_result(
            BrainRouteRequest(
                text=_office_prompt_seed(kind, args),
                member_id=member_id,
                privacy_level="medium",
                estimated_input_tokens=512,
                available_brains=available_brains,
                model_routing_config=config,
                requires_tool_calling=False,
            )
        )
        route = route_selection.route
        if route is None:
            return None
        brain = await self._brains.get_brain(route.primary_brain_id)
        if brain is None:
            return None
        span_id = await self._trace.start_span(
            trace_id,
            span_type=TraceSpanType.MODEL_CALL,
            name="compose office content",
            input_data={
                "office_kind": kind,
                "task_id": task_id,
                "input_summary": str(
                    redact(args.get("goal") or args.get("title") or args.get("content") or "")
                )[:240],
            },
            metadata={
                "brain_id": brain["brain_id"],
                "provider": brain.get("provider"),
                "model_name": brain.get("model_name"),
                "purpose": "office_content_composer",
            },
        )
        if self._audit is not None and not brain.get("is_local"):
            await self._audit.write_event(
                actor_type="system",
                actor_id=member_id,
                action="model_call.cloud_used",
                object_type="brain",
                object_id=brain["brain_id"],
                summary="Office 内容润色使用了云端模型",
                risk_level=RiskLevel.R2,
                payload={"brain_id": brain["brain_id"], "task_id": task_id},
                trace_id=trace_id,
            )
        client = OpenAICompatibleClient(
            str(brain["endpoint"]),
            self._secrets.get_secret(brain.get("api_key_ref")),
        )
        try:
            result = await client.complete_chat(
                ModelChatRequest(
                    model=str(brain["model_name"]),
                    messages=_office_model_messages(kind, args, fallback),
                    temperature=0.2,
                    max_output_tokens=384,
                    top_p=0.9,
                    timeout_seconds=min(int(brain.get("timeout_seconds") or 180), 30),
                    stream=False,
                    trace_id=trace_id,
                    turn_id=f"task:{task_id or 'office'}",
                    route_id=f"office:{brain['brain_id']}",
                    privacy_level="medium",
                    retry_count=0,
                    metadata={"purpose": "office_content_composer"},
                ),
                CancelToken(),
            )
            parsed = _parse_office_model_result(result.text)
        except (ModelAdapterError, json.JSONDecodeError, TypeError, ValueError) as exc:
            await self._trace.end_span(
                span_id,
                status=TraceSpanStatus.FAILED,
                output_data={"error": str(redact(str(exc)))[:240]},
            )
            return None
        await self._trace.end_span(
            span_id,
            output_data={
                "finish_reason": result.finish_reason,
                "usage": result.usage,
                "model_used": True,
            },
        )
        return OfficeComposedContent(
            title=_text(parsed.get("title"), fallback.title),
            summary=_text(parsed.get("summary"), fallback.summary),
            model_used=True,
            strategy="model_routing",
            brain_id=str(brain["brain_id"]),
            finish_reason=result.finish_reason,
        )


class OfficeToolService:
    def __init__(
        self,
        artifact_store: ArtifactStore,
        *,
        brain_repo: BrainRepository | None = None,
        model_routing_service: ModelRoutingService | None = None,
        secret_store: SecretStore | None = None,
        trace_service: TraceService | None = None,
        audit_service: AuditEventService | None = None,
    ) -> None:
        self._artifacts = artifact_store
        self._composer = OfficeContentComposer(
            brain_repo=brain_repo,
            model_routing_service=model_routing_service,
            secret_store=secret_store,
            trace_service=trace_service,
            audit_service=audit_service,
        )

    async def execute(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        if not request.task_id:
            raise AppError(
                ErrorCode.TOOL_PERMISSION_DENIED,
                "Office 工具必须绑定任务",
                status_code=422,
            )
        self._require_supported_provider(request)
        if request.tool_name == "office.word.generate":
            return await self._generate_word(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.word.edit":
            return await self._edit_word(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.excel.generate":
            return await self._generate_excel(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.excel.edit":
            return await self._edit_excel(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.ppt.generate":
            return await self._generate_ppt(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.ppt.edit":
            return await self._edit_ppt(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.mail.draft":
            return await self._mail_draft(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.mail.send":
            return await self._mail_send(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name == "office.calendar.plan":
            return await self._calendar_plan(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        if request.tool_name in {
            "office.document.share",
            "office.document.delete",
            "office.document.overwrite",
            "office.document.modify_shared",
        }:
            return await self._governed_document_action(
                request,
                tool_call_id=tool_call_id,
                organization_id=organization_id,
                trace_id=trace_id,
            )
        raise AppError(ErrorCode.TOOL_NOT_FOUND, "Office 工具不存在", status_code=404)

    def _require_supported_provider(self, request: ToolExecuteRequest) -> None:
        provider_ref = str(request.args.get("provider_ref") or LOCAL_OFFICE_PROVIDER.provider_ref)
        if provider_ref != LOCAL_OFFICE_PROVIDER.provider_ref:
            raise AppError(
                ErrorCode.TOOL_SCHEMA_INVALID,
                "Office provider is unavailable or not authorized",
                status_code=422,
                details={
                    "provider_ref": provider_ref,
                    "blocked_reason": "provider_unavailable_or_not_authorized",
                },
            )

    async def _generate_word(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        _require_docx()
        args = request.args
        composed = await self._composer.compose(
            "word",
            args,
            trace_id=trace_id,
            task_id=request.task_id,
            member_id=request.member_id,
        )
        title = composed.title or "Office Word Document"
        filename = _ensure_suffix(str(args.get("filename") or "office-document.docx"), ".docx")
        document = Document()
        _set_docx_defaults(document)
        document.add_heading(title, level=0)
        summary = composed.summary
        if summary:
            document.add_paragraph(summary)
        for section in _sections(args, default_title="正文"):
            document.add_heading(section["title"], level=1)
            for paragraph in section["paragraphs"]:
                document.add_paragraph(paragraph)
            for bullets in section["bullets"]:
                document.add_paragraph(bullets, style="List Bullet")
        for table in _tables(args):
            _append_docx_table(document, table)
        docx_content = _save_docx(document)
        artifact = await self._artifacts.write_bytes(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            tool_call_id=tool_call_id,
            display_name=filename,
            content=docx_content,
            artifact_type="office_document",
            content_type=DOCX_CONTENT_TYPE,
            metadata=_office_metadata("word", "generate", args, composed),
            trace_id=trace_id,
        )
        return _office_result("word_document", artifact, model_used=False), [artifact]

    async def _edit_word(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        _require_docx()
        source = await self._source_artifact(request)
        document = Document(str(self._artifacts.path_for_artifact(source)))
        for operation in _operations(request.args):
            op_type = str(operation.get("type") or operation.get("op") or "").strip()
            if op_type in {"append_section", "add_section"}:
                document.add_heading(_text(operation.get("title"), "新增章节"), level=1)
                content = operation.get("paragraphs") or operation.get("content")
                for paragraph in _as_text_list(content):
                    document.add_paragraph(paragraph)
                for bullet in _as_text_list(operation.get("bullets")):
                    document.add_paragraph(bullet, style="List Bullet")
            elif op_type == "replace_text":
                _replace_docx_text(
                    document,
                    str(operation.get("old") or ""),
                    str(operation.get("new") or ""),
                )
            elif op_type == "add_table":
                _append_docx_table(document, operation)
            else:
                raise AppError(
                    ErrorCode.TOOL_SCHEMA_INVALID,
                    "不支持的 Word 编辑操作",
                    status_code=422,
                    details={"operation": op_type},
                )
        filename = _ensure_suffix(
            str(request.args.get("filename") or _versioned_name(source.display_name, "edited")),
            ".docx",
        )
        artifact = await self._artifacts.write_bytes(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            tool_call_id=tool_call_id,
            display_name=filename,
            content=_save_docx(document),
            artifact_type="office_document",
            content_type=DOCX_CONTENT_TYPE,
            metadata={
                **_office_metadata(
                    "word",
                    "edit",
                    request.args,
                    await self._composer.compose(
                        "word",
                        request.args,
                        trace_id=trace_id,
                        task_id=request.task_id,
                        member_id=request.member_id,
                    ),
                ),
                "source_artifact_id": source.artifact_id,
                "source_checksum": source.checksum,
            },
            trace_id=trace_id,
        )
        return _office_result("word_document", artifact, source=source), [artifact]

    async def _generate_excel(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        _require_xlsx()
        args = request.args
        composed = await self._composer.compose(
            "excel",
            args,
            trace_id=trace_id,
            task_id=request.task_id,
            member_id=request.member_id,
        )
        workbook = Workbook()
        default_sheet = workbook.active
        sheets = _sheets(args)
        for index, sheet in enumerate(sheets):
            if index == 0 and args.get("summary") and not sheet.get("summary"):
                sheet = {**sheet, "summary": args["summary"]}
            worksheet = default_sheet if index == 0 else workbook.create_sheet()
            worksheet.title = _sheet_title(
                sheet.get("name") or sheet.get("title") or f"Sheet{index + 1}"
            )
            _fill_sheet(worksheet, sheet)
        if len(workbook.sheetnames) > len(sheets):
            workbook.remove(default_sheet)
        _maybe_add_summary_sheet(workbook, args)
        filename = _ensure_suffix(str(args.get("filename") or "office-workbook.xlsx"), ".xlsx")
        artifact = await self._artifacts.write_bytes(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            tool_call_id=tool_call_id,
            display_name=filename,
            content=_save_xlsx(workbook),
            artifact_type="office_spreadsheet",
            content_type=XLSX_CONTENT_TYPE,
            metadata=_office_metadata("excel", "generate", args, composed),
            trace_id=trace_id,
        )
        return _office_result("excel_workbook", artifact, model_used=False), [artifact]

    async def _edit_excel(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        _require_xlsx()
        source = await self._source_artifact(request)
        workbook = load_workbook(self._artifacts.path_for_artifact(source))
        for operation in _operations(request.args):
            op_type = str(operation.get("type") or operation.get("op") or "").strip()
            if op_type in {"add_sheet", "append_sheet"}:
                worksheet = workbook.create_sheet(
                    _sheet_title(operation.get("name") or "New Sheet")
                )
                _fill_sheet(worksheet, operation)
            elif op_type == "append_rows":
                worksheet = _worksheet(workbook, operation.get("sheet"))
                for row in operation.get("rows") or []:
                    worksheet.append(_row_values(row, _headers_from_sheet(worksheet)))
            elif op_type == "set_cell":
                worksheet = _worksheet(workbook, operation.get("sheet"))
                worksheet[str(operation.get("cell") or "A1")] = operation.get("value")
            else:
                raise AppError(
                    ErrorCode.TOOL_SCHEMA_INVALID,
                    "不支持的 Excel 编辑操作",
                    status_code=422,
                    details={"operation": op_type},
                )
        filename = _ensure_suffix(
            str(request.args.get("filename") or _versioned_name(source.display_name, "edited")),
            ".xlsx",
        )
        artifact = await self._artifacts.write_bytes(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            tool_call_id=tool_call_id,
            display_name=filename,
            content=_save_xlsx(workbook),
            artifact_type="office_spreadsheet",
            content_type=XLSX_CONTENT_TYPE,
            metadata={
                **_office_metadata(
                    "excel",
                    "edit",
                    request.args,
                    await self._composer.compose(
                        "excel",
                        request.args,
                        trace_id=trace_id,
                        task_id=request.task_id,
                        member_id=request.member_id,
                    ),
                ),
                "source_artifact_id": source.artifact_id,
                "source_checksum": source.checksum,
            },
            trace_id=trace_id,
        )
        return _office_result("excel_workbook", artifact, source=source), [artifact]

    async def _generate_ppt(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        _require_pptx()
        args = request.args
        composed = await self._composer.compose(
            "ppt",
            args,
            trace_id=trace_id,
            task_id=request.task_id,
            member_id=request.member_id,
        )
        presentation = Presentation()
        title = composed.title or "Office Briefing"
        subtitle = composed.summary
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        slides = _slides(args)
        if not slides:
            slides = [{"title": "概览", "bullets": _as_text_list(args.get("content") or title)}]
        for item in slides:
            _append_ppt_slide(presentation, item)
        filename = _ensure_suffix(str(args.get("filename") or "office-briefing.pptx"), ".pptx")
        artifact = await self._artifacts.write_bytes(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            tool_call_id=tool_call_id,
            display_name=filename,
            content=_save_pptx(presentation),
            artifact_type="office_presentation",
            content_type=PPTX_CONTENT_TYPE,
            metadata=_office_metadata("ppt", "generate", args, composed),
            trace_id=trace_id,
        )
        return _office_result("ppt_presentation", artifact, model_used=False), [artifact]

    async def _edit_ppt(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        _require_pptx()
        source = await self._source_artifact(request)
        presentation = Presentation(str(self._artifacts.path_for_artifact(source)))
        for operation in _operations(request.args):
            op_type = str(operation.get("type") or operation.get("op") or "").strip()
            if op_type in {"append_slide", "add_slide"}:
                _append_ppt_slide(presentation, operation)
            elif op_type == "replace_text":
                _replace_ppt_text(
                    presentation,
                    str(operation.get("old") or ""),
                    str(operation.get("new") or ""),
                )
            else:
                raise AppError(
                    ErrorCode.TOOL_SCHEMA_INVALID,
                    "不支持的 PPT 编辑操作",
                    status_code=422,
                    details={"operation": op_type},
                )
        filename = _ensure_suffix(
            str(request.args.get("filename") or _versioned_name(source.display_name, "edited")),
            ".pptx",
        )
        artifact = await self._artifacts.write_bytes(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            tool_call_id=tool_call_id,
            display_name=filename,
            content=_save_pptx(presentation),
            artifact_type="office_presentation",
            content_type=PPTX_CONTENT_TYPE,
            metadata={
                **_office_metadata(
                    "ppt",
                    "edit",
                    request.args,
                    await self._composer.compose(
                        "ppt",
                        request.args,
                        trace_id=trace_id,
                        task_id=request.task_id,
                        member_id=request.member_id,
                    ),
                ),
                "source_artifact_id": source.artifact_id,
                "source_checksum": source.checksum,
            },
            trace_id=trace_id,
        )
        return _office_result("ppt_presentation", artifact, source=source), [artifact]

    async def _mail_draft(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        del tool_call_id
        args = request.args
        subject = str(args.get("title") or "mail-draft")
        recipients = [str(item) for item in list(args.get("recipients") or []) if str(item).strip()]
        body = str(args.get("content") or args.get("summary") or "")
        content = "\n".join(
            [
                f"Subject: {subject}",
                f"To: {', '.join(recipients) if recipients else '(none)'}",
                "",
                body,
            ]
        ).strip()
        artifact = await self._artifacts.write_text(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            display_name=_ensure_suffix(str(args.get("filename") or "mail-draft.md"), ".md"),
            content=content,
            artifact_type="mail_draft",
            metadata={
                "office_kind": "mail",
                "office_operation": "draft",
                "recipients": recipients,
                "provider_ref": str(args.get("provider_ref") or "local.office_suite"),
            },
            trace_id=trace_id,
        )
        return {
            **office_governance_result(tool_name=request.tool_name, args=args, artifact=artifact),
            "subject": subject,
            "recipients": recipients,
            "draft_artifact": artifact.model_dump(mode="json"),
        }, [artifact]

    async def _mail_send(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        del tool_call_id
        args = request.args
        artifact = await self._artifacts.write_text(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            display_name=_ensure_suffix(str(args.get("filename") or "mail-send-record.md"), ".md"),
            content="\n".join(
                [
                    f"Action: {request.tool_name}",
                    f"Subject: {str(args.get('title') or 'mail-send')}",
                    f"To: {', '.join(str(item) for item in list(args.get('recipients') or []))}",
                    "",
                    str(args.get("content") or args.get("summary") or ""),
                ]
            ).strip(),
            artifact_type="mail_send_record",
            metadata={
                "office_kind": "mail",
                "office_operation": "send",
                "provider_ref": str(args.get("provider_ref") or "local.office_suite"),
                "high_risk_actions": list(args.get("high_risk_actions") or []),
            },
            trace_id=trace_id,
        )
        return {
            **office_governance_result(tool_name=request.tool_name, args=args, artifact=artifact),
            "subject": str(args.get("title") or "mail-send"),
            "recipients": [str(item) for item in list(args.get("recipients") or []) if str(item).strip()],
        }, [artifact]

    async def _calendar_plan(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        del tool_call_id
        args = request.args
        title = str(args.get("title") or "calendar-plan")
        attendees = [str(item) for item in list(args.get("attendees") or []) if str(item).strip()]
        scheduled_time = str(args.get("scheduled_time") or "")
        content = "\n".join(
            [
                f"Title: {title}",
                f"When: {scheduled_time or '(unscheduled)'}",
                f"Attendees: {', '.join(attendees) if attendees else '(none)'}",
                "",
                str(args.get("summary") or args.get("content") or ""),
            ]
        ).strip()
        artifact = await self._artifacts.write_text(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            display_name=_ensure_suffix(str(args.get("filename") or "calendar-plan.md"), ".md"),
            content=content,
            artifact_type="calendar_plan",
            metadata={
                "office_kind": "calendar",
                "office_operation": "plan",
                "attendees": attendees,
                "scheduled_time": scheduled_time,
                "provider_ref": str(args.get("provider_ref") or "local.office_suite"),
            },
            trace_id=trace_id,
        )
        return {
            **office_governance_result(tool_name=request.tool_name, args=args, artifact=artifact),
            "title": title,
            "attendees": attendees,
            "scheduled_time": scheduled_time or None,
            "calendar_artifact": artifact.model_dump(mode="json"),
        }, [artifact]

    async def _governed_document_action(
        self,
        request: ToolExecuteRequest,
        *,
        tool_call_id: str,
        organization_id: str,
        trace_id: str | None,
    ) -> tuple[dict[str, Any], list[TaskArtifact]]:
        del tool_call_id
        args = request.args
        title = str(args.get("title") or request.tool_name)
        artifact = await self._artifacts.write_text(
            task_id=request.task_id or "",
            organization_id=organization_id,
            step_id=request.step_id,
            display_name=_ensure_suffix(str(args.get("filename") or "office-governed-action.md"), ".md"),
            content="\n".join(
                [
                    f"Action: {request.tool_name}",
                    f"Title: {title}",
                    f"Source Artifact: {str(args.get('source_artifact_id') or '(none)')}",
                    f"Share Targets: {', '.join(str(item) for item in list(args.get('share_targets') or []))}",
                    "",
                    str(args.get("summary") or args.get("content") or ""),
                ]
            ).strip(),
            artifact_type="office_action_record",
            metadata={
                "office_kind": str(args.get("request_type") or "document"),
                "office_operation": str(args.get("operation") or request.tool_name),
                "provider_ref": str(args.get("provider_ref") or "local.office_suite"),
                "high_risk_actions": list(args.get("high_risk_actions") or []),
            },
            trace_id=trace_id,
        )
        return office_governance_result(tool_name=request.tool_name, args=args, artifact=artifact), [artifact]

    async def _source_artifact(self, request: ToolExecuteRequest) -> TaskArtifact:
        artifact_id = str(request.args.get("source_artifact_id") or "").strip()
        if not artifact_id:
            raise AppError(
                ErrorCode.TOOL_SCHEMA_INVALID,
                "source_artifact_id 必填",
                status_code=422,
            )
        artifact, _preview = await self._artifacts.read_preview(artifact_id, limit=32)
        if artifact.task_id != request.task_id:
            raise AppError(
                ErrorCode.TOOL_PERMISSION_DENIED,
                "不能编辑其他任务的 artifact",
                status_code=403,
            )
        return artifact


def _set_docx_defaults(document: DocxDocument) -> None:
    style = document.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(10.5)


def _require_docx() -> None:
    if DOCX_IMPORT_ERROR is not None:
        raise AppError(
            ErrorCode.TOOL_EXECUTION_FAILED,
            "Office Word 工具依赖 python-docx 未安装",
            status_code=500,
        ) from DOCX_IMPORT_ERROR


def _require_xlsx() -> None:
    if XLSX_IMPORT_ERROR is not None:
        raise AppError(
            ErrorCode.TOOL_EXECUTION_FAILED,
            "Office Excel 工具依赖 openpyxl 未安装",
            status_code=500,
        ) from XLSX_IMPORT_ERROR


def _require_pptx() -> None:
    if PPTX_IMPORT_ERROR is not None:
        raise AppError(
            ErrorCode.TOOL_EXECUTION_FAILED,
            "Office PPT 工具依赖 python-pptx 未安装",
            status_code=500,
        ) from PPTX_IMPORT_ERROR


def _append_docx_table(document: DocxDocument, table_spec: dict[str, Any]) -> None:
    headers = [str(item) for item in table_spec.get("headers") or []]
    rows = table_spec.get("rows") or []
    if not headers and rows and isinstance(rows[0], dict):
        headers = [str(key) for key in rows[0].keys()]
    if not headers:
        return
    document.add_paragraph()
    table = document.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for index, header in enumerate(headers):
        cell = table.rows[0].cells[index]
        cell.text = header
        for run in cell.paragraphs[0].runs:
            run.bold = True
    for raw_row in rows:
        cells = table.add_row().cells
        values = _row_values(raw_row, headers)
        for index, value in enumerate(values[: len(headers)]):
            cells[index].text = str(value)


def _replace_docx_text(document: DocxDocument, old: str, new: str) -> None:
    if not old:
        return
    for paragraph in document.paragraphs:
        if old in paragraph.text:
            for run in paragraph.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if old in paragraph.text:
                        for run in paragraph.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)


def _save_docx(document: DocxDocument) -> bytes:
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _sheets(args: dict[str, Any]) -> list[dict[str, Any]]:
    raw = args.get("sheets")
    if isinstance(raw, list) and raw:
        return [item for item in raw if isinstance(item, dict)] or [{"name": "Sheet1"}]
    return [
        {
            "name": args.get("sheet_name") or "Sheet1",
            "headers": args.get("headers") or ["项目", "数值"],
            "rows": args.get("rows")
            or [
                {"项目": "输入", "数值": _text(args.get("content") or args.get("goal"), "demo")},
            ],
            "summary": args.get("summary"),
            "add_totals": args.get("add_totals"),
            "chart": args.get("chart"),
        }
    ]


def _fill_sheet(worksheet: Any, sheet: dict[str, Any]) -> None:
    summary = _text(sheet.get("summary"), "")
    row_offset = 1
    if summary:
        worksheet.cell(row=1, column=1, value=summary)
        worksheet.cell(row=1, column=1).font = Font(bold=True)
        row_offset = 3
    headers = [str(item) for item in sheet.get("headers") or []]
    rows = sheet.get("rows") or []
    if not headers and rows and isinstance(rows[0], dict):
        headers = [str(key) for key in rows[0].keys()]
    if not headers:
        headers = ["项目", "数值"]
    for column, header in enumerate(headers, start=1):
        cell = worksheet.cell(row=row_offset, column=column, value=header)
        cell.font = Font(bold=True)
    for row_index, raw_row in enumerate(rows, start=row_offset + 1):
        for column, value in enumerate(_row_values(raw_row, headers), start=1):
            worksheet.cell(row=row_index, column=column, value=value)
    if sheet.get("add_totals") and rows:
        total_row = row_offset + len(rows) + 1
        worksheet.cell(row=total_row, column=1, value="合计")
        for column in range(2, len(headers) + 1):
            letter = worksheet.cell(row=row_offset, column=column).column_letter
            worksheet.cell(
                row=total_row,
                column=column,
                value=f"=SUM({letter}{row_offset + 1}:{letter}{total_row - 1})",
            )
    if sheet.get("chart") and rows and len(headers) >= 2:
        chart = BarChart()
        chart.title = str(sheet.get("chart_title") or "数据对比")
        data = Reference(
            worksheet,
            min_col=2,
            min_row=row_offset,
            max_row=row_offset + len(rows),
            max_col=min(len(headers), 3),
        )
        cats = Reference(
            worksheet,
            min_col=1,
            min_row=row_offset + 1,
            max_row=row_offset + len(rows),
        )
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        worksheet.add_chart(chart, "E2")
    for column_cells in worksheet.columns:
        max_length = max(len(str(cell.value or "")) for cell in column_cells)
        worksheet.column_dimensions[column_cells[0].column_letter].width = min(
            max(max_length + 2, 12),
            42,
        )


def _maybe_add_summary_sheet(workbook: OpenpyxlWorkbook, args: dict[str, Any]) -> None:
    summary = _text(args.get("analysis_summary") or args.get("summary"), "")
    if not summary:
        return
    worksheet = workbook.create_sheet("Summary")
    worksheet["A1"] = "分析摘要"
    worksheet["A1"].font = Font(bold=True)
    worksheet["A2"] = summary
    worksheet.column_dimensions["A"].width = 60


def _worksheet(workbook: OpenpyxlWorkbook, name: Any) -> Any:
    if name and str(name) in workbook.sheetnames:
        return workbook[str(name)]
    return workbook.active


def _headers_from_sheet(worksheet: Any) -> list[str]:
    return [str(cell.value) for cell in next(worksheet.iter_rows(min_row=1, max_row=1))]


def _save_xlsx(workbook: OpenpyxlWorkbook) -> bytes:
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _append_ppt_slide(presentation: PptxPresentation, spec: dict[str, Any]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = _text(spec.get("title"), "新页面")
    body = slide.placeholders[1].text_frame
    body.clear()
    bullets = _as_text_list(spec.get("bullets") or spec.get("content"))
    if not bullets:
        bullets = ["待补充"]
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        for run in paragraph.runs:
            run.font.size = PptPt(20)
    notes = _text(spec.get("speaker_notes"), "")
    if notes:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = notes
    if spec.get("footer"):
        left = Inches(0.6)
        top = Inches(6.9)
        width = Inches(8.8)
        height = Inches(0.3)
        box = slide.shapes.add_textbox(left, top, width, height)
        box.text_frame.text = _text(spec.get("footer"), "")
        for paragraph in box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptPt(9)


def _replace_ppt_text(presentation: PptxPresentation, old: str, new: str) -> None:
    if not old:
        return
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def _save_pptx(presentation: PptxPresentation) -> bytes:
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _sections(args: dict[str, Any], *, default_title: str) -> list[dict[str, Any]]:
    raw = args.get("sections")
    if isinstance(raw, list) and raw:
        result = []
        for item in raw:
            if isinstance(item, dict):
                result.append(
                    {
                        "title": _text(item.get("title"), default_title),
                        "paragraphs": _as_text_list(item.get("paragraphs") or item.get("content")),
                        "bullets": [_text(bullet, "") for bullet in item.get("bullets") or []],
                    }
                )
            else:
                result.append(
                    {
                        "title": default_title,
                        "paragraphs": [_text(item, "")],
                        "bullets": [],
                    }
                )
        return result
    content = _text(args.get("content") or args.get("goal"), "")
    return [{"title": default_title, "paragraphs": [content] if content else [], "bullets": []}]


def _tables(args: dict[str, Any]) -> list[dict[str, Any]]:
    raw = args.get("tables")
    if isinstance(raw, list):
        return [item for item in raw if isinstance(item, dict)]
    table = args.get("table")
    return [table] if isinstance(table, dict) else []


def _slides(args: dict[str, Any]) -> list[dict[str, Any]]:
    raw = args.get("slides")
    if isinstance(raw, list):
        return [item for item in raw if isinstance(item, dict)]
    return []


def _operations(args: dict[str, Any]) -> list[dict[str, Any]]:
    operations = args.get("operations")
    if isinstance(operations, list) and operations:
        return [item for item in operations if isinstance(item, dict)]
    operation = args.get("operation")
    if isinstance(operation, dict):
        return [operation]
    raise AppError(ErrorCode.TOOL_SCHEMA_INVALID, "operations 必填", status_code=422)


def _row_values(row: Any, headers: list[str]) -> list[Any]:
    if isinstance(row, dict):
        return [_coerce_cell_value(row.get(header)) for header in headers]
    if isinstance(row, (list, tuple)):
        return [_coerce_cell_value(item) for item in row]
    return [_coerce_cell_value(row)]


def _coerce_cell_value(value: Any) -> Any:
    if value is None:
        return ""
    if isinstance(value, int | float | bool):
        return value
    return str(value)


def _as_text_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        return [_text(item, "") for item in value if _text(item, "")]
    text = _text(value, "")
    if not text:
        return []
    lines = [line.strip(" -\t") for line in text.splitlines() if line.strip()]
    return lines or [text]


def _text(value: Any, default: str) -> str:
    if value is None:
        return default
    text = str(redact(value)).strip()
    return text or default


def _default_title(args: dict[str, Any]) -> str:
    content = _text(args.get("goal") or args.get("content"), "")
    if not content:
        return "Office Artifact"
    return content[:48]


def _fallback_title(kind: str) -> str:
    if kind == "word":
        return "Office Word Document"
    if kind == "excel":
        return "Office Workbook"
    if kind == "ppt":
        return "Office Briefing"
    return "Office Artifact"


def _ensure_suffix(filename: str, suffix: str) -> str:
    clean = re.sub(r"[\\/:*?\"<>|]+", "_", filename.strip()) or f"artifact{suffix}"
    return clean if clean.lower().endswith(suffix) else f"{clean}{suffix}"


def _versioned_name(name: str, marker: str) -> str:
    if "." not in name:
        return f"{name}-{marker}"
    stem, suffix = name.rsplit(".", 1)
    return f"{stem}-{marker}.{suffix}"


def _sheet_title(value: Any) -> str:
    title = re.sub(r"[:\\/?*\\[\\]]+", "_", _text(value, "Sheet1"))[:31]
    return title or "Sheet1"


def _office_prompt_seed(kind: str, args: dict[str, Any]) -> str:
    raw = args.get("goal") or args.get("title") or args.get("content") or kind
    return str(redact(raw))[:1000]


def _office_model_messages(
    kind: str,
    args: dict[str, Any],
    fallback: OfficeComposedContent,
) -> list[dict[str, str]]:
    payload = {
        "office_kind": kind,
        "goal": str(redact(args.get("goal") or ""))[:1000],
        "content": str(redact(args.get("content") or args.get("summary") or ""))[:2000],
        "fallback": {
            "title": fallback.title,
            "summary": fallback.summary,
        },
    }
    return [
        {
            "role": "system",
            "content": (
                "Return only compact JSON with keys title and summary. "
                "Polish office document metadata only; do not invent private facts."
            ),
        },
        {"role": "user", "content": json.dumps(payload, ensure_ascii=False)},
    ]


def _parse_office_model_result(raw: str) -> dict[str, Any]:
    text = raw.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    payload = json.loads(text)
    if not isinstance(payload, dict):
        raise TypeError("office composer model output must be an object")
    return payload


def _office_metadata(
    kind: str,
    operation: str,
    args: dict[str, Any],
    composed: OfficeComposedContent,
) -> dict[str, Any]:
    input_summary = args.get("goal") or args.get("title") or args.get("content") or ""
    return {
        "office_kind": kind,
        "office_operation": operation,
        "model_used": composed.model_used,
        "model_fallback": None if composed.model_used else composed.strategy,
        "composer_strategy": composed.strategy,
        "composer_brain_id": composed.brain_id,
        "composer_finish_reason": composed.finish_reason,
        "input_summary": str(redact(input_summary))[:240],
        "redaction_summary": {"policy": "trace_service.redact"},
    }


def _office_result(
    key: str,
    artifact: TaskArtifact,
    *,
    source: TaskArtifact | None = None,
    model_used: bool = False,
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        key: artifact.model_dump(mode="json"),
        "artifact_id": artifact.artifact_id,
        "uri": artifact.uri,
        "content_type": artifact.content_type,
        "checksum": artifact.checksum,
        "model_used": model_used,
    }
    if source is not None:
        payload["source_artifact_id"] = source.artifact_id
    return payload
