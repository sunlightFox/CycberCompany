from __future__ import annotations

import importlib.util
from typing import Any, cast

import pytest
from app.services.registry import ServiceRegistry
from fastapi.testclient import TestClient

OFFICE_DEPS_AVAILABLE = all(
    importlib.util.find_spec(name) is not None for name in ["docx", "openpyxl", "pptx"]
)


def test_phase56_office_tools_are_registered(client: TestClient) -> None:
    tools = client.get("/api/tools").json()["items"]
    names = {item["tool_name"] for item in tools}
    assert {
        "office.word.generate",
        "office.word.edit",
        "office.excel.generate",
        "office.excel.edit",
        "office.ppt.generate",
        "office.ppt.edit",
    }.issubset(names)


def test_phase56_clawhub_office_catalog_has_word_excel_ppt(client: TestClient) -> None:
    for query, expected in [
        ("word", "official/office/word-report"),
        ("excel", "official/office/excel-analysis-workbook"),
        ("ppt", "official/office/ppt-briefing"),
    ]:
        response = client.get(
            "/api/skills/catalog/search",
            params={"q": query, "repository_id": "clawhub", "limit": 20},
        )
        assert response.status_code == 200, response.text
        refs = {item["package_ref"] for item in response.json()["items"]}
        assert expected in refs


def test_phase56_office_tool_requires_task(client: TestClient) -> None:
    response = client.post(
        "/api/tools/execute",
        json={
            "tool_name": "office.word.generate",
            "args": {"title": "无任务文档"},
        },
    )
    assert response.status_code == 422
    assert response.json()["error"]["code"] == "TOOL_PERMISSION_DENIED"


@pytest.mark.skipif(
    not OFFICE_DEPS_AVAILABLE,
    reason="office python dependencies are not installed",
)
def test_phase56_word_excel_ppt_generate_and_edit(client: TestClient) -> None:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation

    app = cast(Any, client.app)
    registry = cast(ServiceRegistry, app.state.registry)
    task = client.post("/api/tasks", json={"goal": "Office 工具测试", "auto_start": False}).json()
    task_id = task["task_id"]

    word = _execute(
        client,
        task_id,
        "office.word.generate",
        {
            "filename": "project-report.docx",
            "title": "项目周报",
            "summary": "本周完成需求评审",
            "sections": [{"title": "进展", "paragraphs": ["完成接口设计"]}],
            "tables": [{"headers": ["事项", "状态"], "rows": [["需求评审", "完成"]]}],
        },
    )
    word_artifact = word["artifacts"][0]
    word_path = registry.artifact_store.path_for_artifact(_artifact(word_artifact))
    doc = Document(str(word_path))
    assert "项目周报" in "\n".join(paragraph.text for paragraph in doc.paragraphs)
    assert word_artifact["metadata"]["model_used"] is False
    assert word_artifact["metadata"]["composer_strategy"] == "deterministic_template"

    edited_word = _execute(
        client,
        task_id,
        "office.word.edit",
        {
            "source_artifact_id": word_artifact["artifact_id"],
            "operations": [
                {"type": "append_section", "title": "风险", "paragraphs": ["资源排期需确认"]}
            ],
        },
    )
    assert edited_word["artifacts"][0]["artifact_id"] != word_artifact["artifact_id"]

    excel = _execute(
        client,
        task_id,
        "office.excel.generate",
        {
            "filename": "sales.xlsx",
            "summary": "销售数据",
            "sheets": [
                {
                    "name": "Data",
                    "headers": ["指标", "数值"],
                    "rows": [["收入", 120], ["成本", 80]],
                    "add_totals": True,
                }
            ],
        },
    )
    excel_artifact = excel["artifacts"][0]
    workbook = load_workbook(registry.artifact_store.path_for_artifact(_artifact(excel_artifact)))
    assert "Data" in workbook.sheetnames
    assert workbook["Data"]["A1"].value == "销售数据"

    edited_excel = _execute(
        client,
        task_id,
        "office.excel.edit",
        {
            "source_artifact_id": excel_artifact["artifact_id"],
            "operations": [{"type": "append_rows", "sheet": "Data", "rows": [["利润", 40]]}],
        },
    )
    assert edited_excel["artifacts"][0]["artifact_id"] != excel_artifact["artifact_id"]

    ppt = _execute(
        client,
        task_id,
        "office.ppt.generate",
        {
            "filename": "briefing.pptx",
            "title": "项目汇报",
            "subtitle": "验收阶段",
            "slides": [
                {"title": "背景", "bullets": ["目标明确"]},
                {"title": "进展", "bullets": ["接口完成"]},
            ],
        },
    )
    ppt_artifact = ppt["artifacts"][0]
    presentation = Presentation(
        str(registry.artifact_store.path_for_artifact(_artifact(ppt_artifact)))
    )
    assert len(presentation.slides) == 3

    edited_ppt = _execute(
        client,
        task_id,
        "office.ppt.edit",
        {
            "source_artifact_id": ppt_artifact["artifact_id"],
            "operations": [{"type": "append_slide", "title": "风险", "bullets": ["排期需确认"]}],
        },
    )
    assert edited_ppt["artifacts"][0]["artifact_id"] != ppt_artifact["artifact_id"]


@pytest.mark.skipif(
    not OFFICE_DEPS_AVAILABLE,
    reason="office python dependencies are not installed",
)
def test_phase56_cross_task_office_edit_is_denied(client: TestClient) -> None:
    task_a = client.post("/api/tasks", json={"goal": "Office A", "auto_start": False}).json()
    task_b = client.post("/api/tasks", json={"goal": "Office B", "auto_start": False}).json()
    generated = _execute(
        client,
        task_a["task_id"],
        "office.word.generate",
        {"filename": "a.docx", "title": "A"},
    )
    response = client.post(
        "/api/tools/execute",
        json={
            "task_id": task_b["task_id"],
            "tool_name": "office.word.edit",
            "args": {
                "source_artifact_id": generated["artifacts"][0]["artifact_id"],
                "operations": [{"type": "append_section", "title": "X"}],
            },
        },
    )
    assert response.status_code == 403
    assert response.json()["error"]["code"] == "TOOL_PERMISSION_DENIED"


def test_phase56_repository_install_enable_grant_and_office_plan(client: TestClient) -> None:
    install = client.post(
        "/api/skills/install",
        json={
            "source_type": "repository_ref",
            "source_uri": "clawhub:official/office/word-report",
        },
    )
    assert install.status_code == 200, install.text
    payload = install.json()
    bundle_id = payload["bundle"]["bundle_id"]
    skill_id = payload["skills"][0]["skill_id"]
    enabled = client.post(
        f"/api/plugins/{bundle_id}/enable",
        json={"actor_member_id": "mem_xiaoyao"},
    )
    assert enabled.status_code == 200, enabled.text
    grant = client.post(
        f"/api/skills/{skill_id}/grants",
        json={"allowed_tools": ["office.word.generate"]},
    )
    assert grant.status_code == 200, grant.text

    task = client.post(
        "/api/tasks",
        json={"goal": "帮我生成一份 Word 项目周报", "auto_start": False},
    ).json()
    plan = client.get(f"/api/tasks/{task['task_id']}/plan").json()["plan"]
    assert any(
        step["step_type"] == "skill_run" and step["input"]["skill_id"] == skill_id
        for step in plan["steps"]
    )


@pytest.mark.skipif(
    not OFFICE_DEPS_AVAILABLE,
    reason="office python dependencies are not installed",
)
def test_phase56_word_skill_run_creates_real_docx_artifact(client: TestClient) -> None:
    from docx import Document

    app = cast(Any, client.app)
    registry = cast(ServiceRegistry, app.state.registry)
    install = client.post(
        "/api/skills/install",
        json={
            "source_type": "repository_ref",
            "source_uri": "clawhub:official/office/word-report",
        },
    )
    assert install.status_code == 200, install.text
    payload = install.json()
    bundle_id = payload["bundle"]["bundle_id"]
    skill_id = payload["skills"][0]["skill_id"]
    assert client.post(
        f"/api/plugins/{bundle_id}/enable",
        json={"actor_member_id": "mem_xiaoyao"},
    ).status_code == 200
    assert client.post(
        f"/api/skills/{skill_id}/grants",
        json={"allowed_tools": ["office.word.generate"]},
    ).status_code == 200

    task = client.post(
        "/api/tasks",
        json={"goal": "帮我生成一份 Word 项目周报", "auto_start": False},
    ).json()
    started = client.post(f"/api/tasks/{task['task_id']}/start")
    assert started.status_code == 200, started.text
    detail = started.json()
    assert detail["status"] == "completed"

    replay = client.get(f"/api/tasks/{task['task_id']}/replay").json()
    skill_runs = replay["skill_runs"]
    assert skill_runs[0]["status"] == "completed"
    artifacts = client.get(f"/api/tasks/{task['task_id']}/artifacts").json()["items"]
    docx_artifact = next(item for item in artifacts if item["content_type"].endswith("document"))
    docx_path = registry.artifact_store.path_for_artifact(_artifact(docx_artifact))
    document = Document(str(docx_path))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert "项目周报" in text


@pytest.mark.skipif(
    not OFFICE_DEPS_AVAILABLE,
    reason="office python dependencies are not installed",
)
@pytest.mark.parametrize(
    ("source_uri", "tool_name", "goal", "content_type_suffix", "loader_check"),
    [
        (
            "clawhub:official/office/excel-analysis-workbook",
            "office.excel.generate",
            "把这些销售数据做成 Excel 分析表",
            "spreadsheetml.sheet",
            "xlsx",
        ),
        (
            "clawhub:official/office/ppt-briefing",
            "office.ppt.generate",
            "做一个 5 页 PPT 汇报",
            "presentationml.presentation",
            "pptx",
        ),
    ],
)
def test_phase56_excel_and_ppt_skill_runs_create_real_office_artifacts(
    client: TestClient,
    source_uri: str,
    tool_name: str,
    goal: str,
    content_type_suffix: str,
    loader_check: str,
) -> None:
    from openpyxl import load_workbook
    from pptx import Presentation

    app = cast(Any, client.app)
    registry = cast(ServiceRegistry, app.state.registry)
    install = client.post(
        "/api/skills/install",
        json={"source_type": "repository_ref", "source_uri": source_uri},
    )
    assert install.status_code == 200, install.text
    payload = install.json()
    bundle_id = payload["bundle"]["bundle_id"]
    skill_id = payload["skills"][0]["skill_id"]
    assert client.post(
        f"/api/plugins/{bundle_id}/enable",
        json={"actor_member_id": "mem_xiaoyao"},
    ).status_code == 200
    assert client.post(
        f"/api/skills/{skill_id}/grants",
        json={"allowed_tools": [tool_name]},
    ).status_code == 200

    task = client.post("/api/tasks", json={"goal": goal, "auto_start": False}).json()
    started = client.post(f"/api/tasks/{task['task_id']}/start")
    assert started.status_code == 200, started.text
    assert started.json()["status"] == "completed"
    artifacts = client.get(f"/api/tasks/{task['task_id']}/artifacts").json()["items"]
    office_artifact = next(
        item for item in artifacts if item["content_type"].endswith(content_type_suffix)
    )
    path = registry.artifact_store.path_for_artifact(_artifact(office_artifact))
    if loader_check == "xlsx":
        workbook = load_workbook(path)
        assert "Data" in workbook.sheetnames
    else:
        presentation = Presentation(str(path))
        assert len(presentation.slides) >= 3


def _execute(
    client: TestClient,
    task_id: str,
    tool_name: str,
    args: dict[str, Any],
) -> dict[str, Any]:
    response = client.post(
        "/api/tools/execute",
        json={"task_id": task_id, "tool_name": tool_name, "args": args},
    )
    assert response.status_code == 200, response.text
    return response.json()


def _artifact(payload: dict[str, Any]) -> Any:
    from core_types import TaskArtifact

    return TaskArtifact(**payload)
