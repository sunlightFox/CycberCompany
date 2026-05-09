from __future__ import annotations

import importlib.util
import time
from typing import Any, cast

import pytest
from app.services.chat_intent_router import (
    ChatIntentRouter,
    browser_search_query,
    browser_search_requires_citation,
    extract_host_software_name,
    host_filesystem_location,
    host_software_action,
    is_browser_search_request,
    is_desktop_native_request,
    is_explicit_download_request,
    is_host_filesystem_list_request,
    is_webpage_read_request,
    office_skill_input,
    parse_office_chat_request,
    webpage_read_url,
)
from app.services.registry import ServiceRegistry
from fastapi.testclient import TestClient

OFFICE_DEPS_AVAILABLE = all(
    importlib.util.find_spec(name) is not None for name in ["docx", "openpyxl", "pptx"]
)


@pytest.mark.parametrize(
    ("text", "route_type", "document_type"),
    [
        ("Office Skill 安装启用后，帮我生成一份 Word 项目周报。", "office_document", "word"),
        ("把这些销售数据做成 Excel 分析表。", "office_document", "excel"),
        ("安装与授权后做一个 5 页 PPT 汇报。", "office_document", "ppt"),
        ("PPT 权限配置好了以后，生成汇报。", "office_document", "ppt"),
        ("Office Skill 生成 Word 项目报告。", "office_document", "word"),
    ],
)
def test_phase57_router_prioritizes_office_over_install_settings_and_clarification(
    text: str,
    route_type: str,
    document_type: str,
) -> None:
    decision = ChatIntentRouter().decide(text)
    assert decision.route_type == route_type
    assert decision.office_request is not None
    assert decision.office_request.document_type == document_type


def test_phase57_router_keeps_download_topic_from_real_download() -> None:
    assert ChatIntentRouter().decide("补一下 artifact 下载端点说明，不要真的下载。").route_type != (
        "browser_download"
    )
    assert is_explicit_download_request("帮我下载 http://127.0.0.1:54069/report.csv。") is True


@pytest.mark.parametrize(
    ("text", "forbidden_route"),
    [
        ("讲一下安装包校验机制，不要安装任何软件。", "host_software_install"),
        ("删除风险说明，不要删除文件。", "file_mutation_task"),
        ("这是一份任务报告，不要生成 Word。", "office_document"),
        ("整理一下汇报思路，不要做成 PPT。", "office_document"),
        ("补一下 artifact 下载端点说明，不要真的下载。", "browser_download"),
    ],
)
def test_phase57_router_respects_negative_action_constraints(
    text: str,
    forbidden_route: str,
) -> None:
    decision = ChatIntentRouter().decide(text)
    assert decision.route_type != forbidden_route


def test_phase57_router_does_not_treat_plain_report_as_ppt() -> None:
    decision = ChatIntentRouter().decide("请总结这次测试汇报的重点，不要创建任务。")
    assert decision.route_type != "office_document"


def test_phase57_router_keeps_skill_mcp_concept_direct_when_user_says_no_task() -> None:
    decision = ChatIntentRouter().decide("解释一下 Skill 和 MCP 有什么区别，不要创建任务。")
    assert decision.route_type == "skill_mcp_concept"


def test_phase57_office_skill_input_extracts_excel_rows_and_ppt_slide_count() -> None:
    excel = parse_office_chat_request(
        "把这些销售数据做成 Excel 分析表：1月收入120成本80，2月收入150成本95。"
    )
    assert excel is not None
    sheets = office_skill_input(excel)["sheets"]
    assert sheets[0]["headers"] == ["期间", "收入", "成本", "利润"]
    assert sheets[0]["rows"] == [["1月", 120, 80, 40], ["2月", 150, 95, 55]]

    ppt = parse_office_chat_request("安装与授权后做一个 5 页 PPT 汇报。")
    assert ppt is not None
    slides = office_skill_input(ppt)["slides"]
    assert len(slides) == 4


def test_phase57_router_keeps_true_host_install_positive() -> None:
    decision = ChatIntentRouter().decide("帮我安装 Example Tool 到这台电脑。")
    assert decision.route_type == "host_software_install"
    assert decision.requires_confirmation is True


def test_phase57_router_treats_host_uninstall_as_host_software_action() -> None:
    decision = ChatIntentRouter().decide("帮我卸载 QQ。")
    assert decision.route_type == "host_software_install"
    assert decision.requires_confirmation is True
    assert host_software_action("帮我卸载 QQ。") == "uninstall"
    assert extract_host_software_name("帮我卸载 QQ。") == "QQ"


@pytest.mark.parametrize(
    ("text", "location"),
    [
        ("我桌面有哪些文件", "desktop"),
        ("桌面有什么文件", "desktop"),
        ("下载目录里有哪些文件", "downloads"),
        ("看看文档目录里有什么", "documents"),
    ],
)
def test_phase57_router_detects_readonly_host_filesystem_list(
    text: str,
    location: str,
) -> None:
    decision = ChatIntentRouter().decide(text)
    assert decision.route_type == "host_filesystem_list"
    assert decision.requires_confirmation is False
    assert decision.metadata["location"] == location
    assert is_host_filesystem_list_request(text) is True
    assert host_filesystem_location(text) == location


@pytest.mark.parametrize(
    "text",
    [
        "帮我看一下这网站有什么内容，https://example.com/news/1",
        "这个网页讲什么 https://example.com/post",
        "总结这个链接：https://example.com/a?b=1",
        "https://example.com/page 这个链接主要说什么",
    ],
)
def test_phase57_router_detects_readonly_webpage_read(text: str) -> None:
    decision = ChatIntentRouter().decide(text)
    assert decision.route_type == "browser_read_page"
    assert decision.requires_confirmation is False
    url = decision.metadata["url"]
    assert isinstance(url, str)
    assert url.startswith("https://example.com")
    assert is_webpage_read_request(text) is True
    extracted_url = webpage_read_url(text)
    assert extracted_url is not None
    assert extracted_url.startswith("https://example.com")


def test_phase57_router_detects_browser_search_with_citation() -> None:
    text = "请用浏览器搜索 chat main chain regression，并总结结果，必须说明证据来源。"
    decision = ChatIntentRouter().decide(text)
    assert decision.route_type == "browser_search_with_citation"
    assert decision.requires_confirmation is False
    assert is_browser_search_request(text) is True
    assert browser_search_requires_citation(text) is True
    assert "chat main chain regression" in browser_search_query(text)


def test_phase57_router_detects_desktop_native_request() -> None:
    text = "请帮我操作桌面窗口，把当前桌面上的记事本窗口最小化。"
    decision = ChatIntentRouter().decide(text)
    assert decision.route_type == "desktop_native_request"
    assert decision.metadata["capability_namespace"] == "desktop"
    assert is_desktop_native_request(text) is True


@pytest.mark.parametrize(
    "text",
    [
        "下载这个链接 https://example.com/report.csv",
        "打开 https://example.com/login 并帮我登录",
        "去 https://example.com/form 提交这个表单",
        "帮我点击 https://example.com 里的购买按钮",
    ],
)
def test_phase57_router_keeps_browser_side_effects_out_of_readonly_route(text: str) -> None:
    assert ChatIntentRouter().decide(text).route_type != "browser_read_page"
    assert is_webpage_read_request(text) is False


def test_phase57_brain_decision_marks_host_uninstall_as_task_request(client: TestClient) -> None:
    body = _turn(client, "phase57-uninstall-brain", "帮我卸载 QQ。")
    assert body["status"] == "completed"
    response = client.get(f"/api/chat/turns/{body['turn_id']}/brain-decision")
    assert response.status_code == 200, response.text
    brain = response.json()
    assert brain["intent"]["primary_intent"] == "task_request"
    assert "host_software_install_request" in brain["intent"]["rule_hits"]
    assert "host_software_change" in brain["intent"]["risk_signals"]


def test_phase57_brain_decision_marks_webpage_read_as_direct_tool(client: TestClient) -> None:
    body = _turn(client, "phase57-browser-read-brain", "这个网页讲什么 https://example.com/post")
    assert body["status"] == "completed"
    response = client.get(f"/api/chat/turns/{body['turn_id']}/brain-decision")
    assert response.status_code == 200, response.text
    brain = response.json()
    assert brain["intent"]["primary_intent"] == "browser_read"
    assert brain["intent"]["needs_tool"] is True
    assert brain["intent"]["needs_task"] is False
    assert brain["mode"]["mode"] == "direct"
    assert brain["clarification"]["needed"] is False


@pytest.mark.skipif(
    not OFFICE_DEPS_AVAILABLE,
    reason="office python dependencies are not installed",
)
@pytest.mark.parametrize(
    ("source_uri", "tool_name", "text", "expected_marker"),
    [
        (
            "clawhub:official/office/word-report",
            "office.word.generate",
            "Office Skill 安装启用后，帮我生成一份 Word 项目周报。",
            "wordprocessingml.document",
        ),
        (
            "clawhub:official/office/excel-analysis-workbook",
            "office.excel.generate",
            "把这些销售数据做成 Excel 分析表。",
            "spreadsheetml.sheet",
        ),
        (
            "clawhub:official/office/ppt-briefing",
            "office.ppt.generate",
            "安装与授权后做一个 5 页 PPT 汇报。",
            "presentationml.presentation",
        ),
    ],
)
def test_phase57_chat_office_auto_runs_enabled_granted_skill(
    client: TestClient,
    source_uri: str,
    tool_name: str,
    text: str,
    expected_marker: str,
) -> None:
    _install_enable_grant(client, source_uri, tool_name)
    body = _turn(client, f"phase57-{tool_name}", text)
    assert body["status"] == "completed"
    events = _events(client, body["turn_id"])
    names = [item["event_type"] for item in events]
    reply = _reply_from_events(events)
    payload = _completed_payload(events)["response_plan"]["structured_payload"]

    assert "task.created" in names
    assert "host_install_plan" not in str(payload)
    assert payload["office_route"]["status"] == "completed"
    assert payload["office_route"]["artifacts"]
    assert payload["office_route"]["artifacts"][0]["download_url"].startswith(
        "/api/artifacts/"
    )
    assert "已经" in reply and "完成" in reply
    assert "<minimax:tool_call" not in reply
    assert "trace_id" not in reply and "task_id" not in reply

    task_id = payload["task_status"]["task_id"]
    artifacts = client.get(f"/api/tasks/{task_id}/artifacts").json()["items"]
    assert any(expected_marker in str(item.get("content_type") or "") for item in artifacts)
    artifact_id = payload["office_route"]["artifacts"][0]["artifact_id"]
    downloaded = client.get(f"/api/artifacts/{artifact_id}/download")
    assert downloaded.status_code == 200
    assert expected_marker in str(downloaded.headers.get("content-type", ""))
    assert downloaded.content


def test_phase57_chat_office_missing_skill_does_not_fake_file(client: TestClient) -> None:
    body = _turn(client, "phase57-office-missing", "Office Skill 生成 Word 项目报告。")
    events = _events(client, body["turn_id"])
    names = [item["event_type"] for item in events]
    reply = _reply_from_events(events)
    payload = _completed_payload(events)["response_plan"]["structured_payload"]

    assert "task.created" not in names
    assert payload["office_route"]["missing_reason"] == "missing_enabled_skill"
    assert (
        "没有假装生成文件" in reply
        or "没有生成文件" in reply
        or "没有假装已经生成" in reply
    )


@pytest.mark.skipif(
    not OFFICE_DEPS_AVAILABLE,
    reason="office python dependencies are not installed",
)
def test_phase57_chat_office_multiturn_generates_edit_and_uses_user_data(
    client: TestClient,
) -> None:
    from core_types import TaskArtifact
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation

    app = cast(Any, client.app)
    registry = cast(ServiceRegistry, app.state.registry)
    _install_enable_grant(
        client,
        "clawhub:official/office/word-report",
        "office.word.generate",
    )
    _install_enable_grant(client, "clawhub:official/office/word-edit", "office.word.edit")
    _install_enable_grant(
        client,
        "clawhub:official/office/excel-analysis-workbook",
        "office.excel.generate",
    )
    _install_enable_grant(client, "clawhub:official/office/ppt-briefing", "office.ppt.generate")

    conversation = client.get("/api/chat/conversations").json()["items"][0]
    conversation_id = conversation["conversation_id"]
    word = _turn(
        client,
        "phase57-multiturn",
        (
            "Office Skill 安装启用后，帮我生成一份 Word 项目周报，"
            "内容包括本周完成接口评审、风险是上线窗口紧、下一步要补自动化测试。"
        ),
        conversation_id=conversation_id,
    )
    edited = _turn(
        client,
        "phase57-multiturn",
        "把刚才生成的 Word 增加风险与下一步章节。",
        conversation_id=conversation_id,
    )
    excel = _turn(
        client,
        "phase57-multiturn",
        "把这些销售数据做成 Excel 分析表：1月收入120成本80，2月收入150成本95。",
        conversation_id=conversation_id,
    )
    ppt = _turn(
        client,
        "phase57-multiturn",
        "安装与授权后做一个 5 页 PPT 汇报，主题是 Q2 增长复盘，面向管理层。",
        conversation_id=conversation_id,
    )

    assert word["status"] == "completed"
    assert edited["status"] == "completed"
    word_events = _events(client, word["turn_id"])
    word_payload = _completed_payload(word_events)["response_plan"]["structured_payload"]
    word_task_id = word_payload["task_status"]["task_id"]
    word_artifacts = client.get(f"/api/tasks/{word_task_id}/artifacts").json()["items"]
    word_docx = _latest_office_artifact(word_artifacts, "wordprocessingml.document")
    word_doc = Document(str(registry.artifact_store.path_for_artifact(TaskArtifact(**word_docx))))
    word_text = "\n".join(paragraph.text for paragraph in word_doc.paragraphs)
    assert "接口评审" in word_text
    assert "上线窗口紧" in word_text
    assert "补自动化测试" in word_text
    original_checksum = word_docx["checksum"]

    edited_events = _events(client, edited["turn_id"])
    edited_payload = _completed_payload(edited_events)["response_plan"]["structured_payload"]
    assert edited_payload["office_route"]["status"] == "completed"
    assert edited_payload["office_route"]["artifacts"]
    edited_task_id = edited_payload["task_status"]["task_id"]
    edited_artifacts = client.get(f"/api/tasks/{edited_task_id}/artifacts").json()["items"]
    edited_docx = _latest_office_artifact(edited_artifacts, "wordprocessingml.document")
    assert edited_docx["checksum"] != original_checksum
    edited_doc = Document(
        str(registry.artifact_store.path_for_artifact(TaskArtifact(**edited_docx)))
    )
    assert "风险与下一步" in "\n".join(paragraph.text for paragraph in edited_doc.paragraphs)

    excel_payload = _completed_payload(_events(client, excel["turn_id"]))["response_plan"][
        "structured_payload"
    ]
    excel_task_id = excel_payload["task_status"]["task_id"]
    excel_artifacts = client.get(f"/api/tasks/{excel_task_id}/artifacts").json()["items"]
    xlsx = _latest_office_artifact(excel_artifacts, "spreadsheetml.sheet")
    workbook = load_workbook(registry.artifact_store.path_for_artifact(TaskArtifact(**xlsx)))
    values = [row for row in workbook["Data"].iter_rows(values_only=True)]
    assert ("1月", 120, 80, 40) in values
    assert ("2月", 150, 95, 55) in values

    ppt_payload = _completed_payload(_events(client, ppt["turn_id"]))["response_plan"][
        "structured_payload"
    ]
    ppt_task_id = ppt_payload["task_status"]["task_id"]
    ppt_artifacts = client.get(f"/api/tasks/{ppt_task_id}/artifacts").json()["items"]
    pptx = _latest_office_artifact(ppt_artifacts, "presentationml.presentation")
    presentation = Presentation(
        str(registry.artifact_store.path_for_artifact(TaskArtifact(**pptx)))
    )
    assert len(presentation.slides) == 5
    title_text = presentation.slides[0].shapes.title.text
    assert "Q2 增长复盘" in title_text or "Q2 增长复盘" in _reply_from_events(
        _events(client, ppt["turn_id"])
    )


def test_phase57_direct_routes_do_not_fail_without_model(client: TestClient) -> None:
    concept = _turn(client, "phase57-direct", "解释一下 Skill 和 MCP 有什么区别，不要创建任务。")
    download_topic = _turn(client, "phase57-direct", "补一下 artifact 下载端点说明，不要真的下载。")
    strategy = _turn(
        client,
        "phase57-direct",
        "在测试速度、覆盖率、真实模型成本之间做取舍，给我一个有理由的建议。",
    )

    assert concept["status"] == "completed"
    assert download_topic["status"] == "completed"
    assert strategy["status"] == "completed"
    assert any(
        marker in _reply_from_events(_events(client, concept["turn_id"]))
        for marker in ("Skill", "方法包")
    )
    assert "不会触发真实下载" in _reply_from_events(_events(client, download_topic["turn_id"]))
    assert "真实模型" in _reply_from_events(_events(client, strategy["turn_id"]))


def test_phase57_output_filter_removes_model_tool_markup() -> None:
    from app.services.chat_safety import ChatVisibleOutputFilter

    text, summary = ChatVisibleOutputFilter.filter_text(
        "我来处理 <minimax:tool_call name=\"office.word.generate\">x</minimax:tool_call> 完成。"
    )
    assert "<minimax:tool_call" not in text
    assert "model_tool_xml" in summary["blocked_terms"]


def _install_enable_grant(client: TestClient, source_uri: str, tool_name: str) -> dict[str, Any]:
    installed = client.post(
        "/api/skills/install",
        json={"source_type": "repository_ref", "source_uri": source_uri},
    )
    assert installed.status_code == 200, installed.text
    payload = installed.json()
    bundle_id = payload["bundle"]["bundle_id"]
    skill_id = payload["skills"][0]["skill_id"]
    enabled = client.post(
        f"/api/plugins/{bundle_id}/enable",
        json={"actor_member_id": "mem_xiaoyao"},
    )
    assert enabled.status_code == 200, enabled.text
    grant = client.post(
        f"/api/skills/{skill_id}/grants",
        json={"allowed_tools": [tool_name]},
    )
    assert grant.status_code == 200, grant.text
    return {"bundle_id": bundle_id, "skill_id": skill_id}


def _turn(
    client: TestClient,
    session_id: str,
    text: str,
    *,
    conversation_id: str | None = None,
) -> dict[str, Any]:
    created = client.post(
        "/api/chat/turn",
        json={
            "member_id": "mem_xiaoyao",
            "session_id": session_id,
            **({"conversation_id": conversation_id} if conversation_id else {}),
            "input": {"type": "text", "text": text},
        },
    )
    assert created.status_code == 200, created.text
    turn_id = created.json()["turn_id"]
    deadline = time.time() + 30
    while time.time() < deadline:
        detail = client.get(f"/api/chat/turns/{turn_id}")
        assert detail.status_code == 200, detail.text
        body = detail.json()
        if body["status"] in {"completed", "failed", "cancelled"}:
            return body
        time.sleep(0.05)
    raise AssertionError(f"turn {turn_id} did not finish")


def _events(client: TestClient, turn_id: str) -> list[dict[str, Any]]:
    response = client.get(f"/api/chat/turns/{turn_id}/events")
    assert response.status_code == 200, response.text
    return response.json()["items"]


def _reply_from_events(events: list[dict[str, Any]]) -> str:
    return "".join(
        str(item.get("payload", {}).get("payload", {}).get("text", ""))
        for item in events
        if item["event_type"] == "response.delta"
    )


def _completed_payload(events: list[dict[str, Any]]) -> dict[str, Any]:
    completed = next(item for item in events if item["event_type"] == "response.completed")
    return dict(completed["payload"]["payload"])


def _latest_office_artifact(artifacts: list[dict[str, Any]], marker: str) -> dict[str, Any]:
    matches = [
        item
        for item in artifacts
        if str(item.get("content_type") or "").endswith(marker)
        and not (item.get("metadata") or {}).get("copied_for_office_edit")
    ]
    assert matches
    return matches[-1]
